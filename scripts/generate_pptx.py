"""
Generate the PowerPoint presentation (.pptx) for the
Cameroon Malnutrition Atlas — CEC 420 project.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGURES = os.path.join(BASE, "reports", "figures")
OUT = os.path.join(BASE, "reports", "CEC420_Malnutrition_Atlas_Presentation_v2.pptx")

# ── colour palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0, 51, 102)
BLUE   = RGBColor(31, 117, 181)
LIGHT  = RGBColor(217, 225, 242)
ORANGE = RGBColor(255, 121, 0)
WHITE  = RGBColor(255, 255, 255)
GREY   = RGBColor(89, 89, 89)
GREEN  = RGBColor(0, 128, 0)
RED    = RGBColor(192, 0, 0)
DARK   = RGBColor(31, 31, 31)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, italic=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_image_or_placeholder(slide, filename, left, top, width, height, caption=""):
    path = os.path.join(FIGURES, filename) if filename else None
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        # placeholder box
        box = add_rect(slide, left, top, width, height, fill=LIGHT, line_color=BLUE, line_width=1.5)
        add_text(slide, f"[ IMAGE AREA ]\n{caption}",
                 left + 0.1, top + height / 2 - 0.4, width - 0.2, 0.8,
                 font_size=13, italic=True, color=BLUE, align=PP_ALIGN.CENTER)
    if caption:
        add_text(slide, caption,
                 left, top + height + 0.05, width, 0.25,
                 font_size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def add_bullet_block(slide, items, left, top, width, height,
                     font_size=14, color=DARK, bullet_char="▸ "):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def slide_header(slide, title, subtitle=None):
    # Navy top bar
    add_rect(slide, 0, 0, 13.33, 1.1, fill=NAVY)
    add_text(slide, title, 0.3, 0.08, 12.0, 0.75,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri Light")
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.72, 12.0, 0.35,
                 font_size=15, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)
    # Thin accent line
    add_rect(slide, 0, 1.1, 13.33, 0.04, fill=ORANGE)


def footer(slide, slide_num, total=22):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=NAVY)
    add_text(slide, "Cameroon Malnutrition Atlas  |  CEC 420  |  SEPO PERRY-BRADLEY DINGA  |  CT23A145",
             0.3, 7.22, 11.5, 0.25, font_size=9, color=LIGHT, align=PP_ALIGN.LEFT)
    add_text(slide, f"{slide_num} / {total}", 12.6, 7.22, 0.65, 0.25,
             font_size=9, color=LIGHT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 2.8, 13.33, 2.9, fill=BLUE)

add_text(sl, "UNIVERSITY OF BUEA", 1.0, 0.25, 11.0, 0.5,
         font_size=16, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "Department of Computer Engineering – Software Engineering",
         1.0, 0.72, 11.0, 0.35, font_size=13, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(sl, "CAMEROON MALNUTRITION ATLAS", 0.5, 1.45, 12.33, 0.9,
         font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri Light")
add_text(sl, "A Data Mining Approach to Predicting and Mapping Child Stunting",
         1.0, 2.3, 11.0, 0.5, font_size=18, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(sl, "SEPO PERRY-BRADLEY DINGA  |  CT23A145", 1.0, 3.1, 11.0, 0.4,
         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "CEC 420 – Data Mining  |  2025 / 2026", 1.0, 3.5, 11.0, 0.35,
         font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "Powered by CRISP-DM  •  XGBoost  •  K-Means  •  Spearman Tests  •  Next.js Webapp",
         0.5, 4.1, 12.0, 0.35, font_size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

add_rect(sl, 0, 7.2, 13.33, 0.3, fill=RGBColor(0, 31, 63))
add_text(sl, "June 2026", 0.3, 7.22, 12.5, 0.25,
         font_size=9, color=LIGHT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Presentation Outline")
footer(sl, 2)

sections = [
    ("01", "Introduction & Problem Statement"),
    ("02", "CRISP-DM Methodology Framework"),
    ("03", "Data Sources & Dataset Overview"),
    ("04", "Exploratory Data Analysis"),
    ("05", "Data Preparation & Feature Engineering"),
    ("06", "Regression Modelling — 10 Models Compared"),
    ("07", "Hypothesis Testing — H1 through H6"),
    ("08", "Classification — WHO Risk Bands"),
    ("09", "Clustering Analysis — K-Means"),
    ("10", "Forecasting to 2026 / 2028"),
    ("11", "Hotspot Rankings & Deployment"),
    ("12", "Key Findings & Recommendations"),
]

col_items = [sections[:6], sections[6:]]
for c, col in enumerate(col_items):
    for r, (num, title) in enumerate(col):
        x = 0.4 + c * 6.5
        y = 1.35 + r * 0.84
        add_rect(sl, x, y, 0.5, 0.55, fill=ORANGE)
        add_text(sl, num, x, y + 0.06, 0.5, 0.45,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title, x + 0.58, y + 0.1, 5.5, 0.45,
                 font_size=14, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Introduction", "The Challenge: Child Stunting in Cameroon")
footer(sl, 3)

# KPI boxes
kpis = [
    ("29%", "National Stunting Rate\n(DHS 2018)"),
    ("37.9%", "Far North Region\n(Highest, 2018)"),
    ("15.4%", "Littoral Region\n(Lowest, 2018)"),
    ("22.5pp", "North–South Gap\n(Percentage Points)"),
]
for i, (val, label) in enumerate(kpis):
    x = 0.4 + i * 3.1
    add_rect(sl, x, 1.25, 2.8, 1.5, fill=BLUE if i % 2 == 0 else ORANGE)
    add_text(sl, val, x, 1.35, 2.8, 0.75,
             font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri Light")
    add_text(sl, label, x, 2.0, 2.8, 0.65,
             font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Why does this matter?", 0.4, 3.0, 12.0, 0.35,
         font_size=16, bold=True, color=NAVY)
add_bullet_block(sl, [
    "Stunting in the first 1,000 days causes irreversible cognitive deficits and perpetuates intergenerational poverty.",
    "The national 29% average conceals a 22.5 pp north-south gap — resources must be geographically targeted.",
    "The Ministry of Public Health needs a principled, data-driven method for prioritising districts.",
    "This project builds the 'Cameroon Malnutrition Atlas' — a CRISP-DM pipeline that produces ranked targeting intelligence.",
], 0.4, 3.4, 12.2, 3.4, font_size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CRISP-DM FRAMEWORK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "CRISP-DM Methodology", "Cross-Industry Standard Process for Data Mining")
footer(sl, 4)

phases = [
    ("1", "Business\nUnderstanding", "Define targets, hypotheses H1-H6, success criteria"),
    ("2", "Data\nUnderstanding", "Source DHS/MICS, WHO, HDX datasets — 50 obs × 20 cols"),
    ("3", "Data\nPreparation", "Clean, impute, engineer 4 composite features"),
    ("4", "Modelling", "10 regression + 2 classification + K-Means + linear forecast"),
    ("5", "Evaluation", "5-fold CV, hypothesis tests, business criteria check"),
    ("6", "Deployment", "Hotspot CSV + best_model.joblib + Next.js web app"),
]
colors = [NAVY, BLUE, RGBColor(0, 112, 192), ORANGE, RGBColor(192, 80, 77), GREEN]

for i, (num, phase, detail) in enumerate(phases):
    x = 0.35 + (i % 3) * 4.2
    y = 1.4 + (i // 3) * 2.7
    add_rect(sl, x, y, 3.7, 0.6, fill=colors[i])
    add_text(sl, f"Phase {num}: {phase}", x, y + 0.05, 3.7, 0.55,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x, y + 0.6, 3.7, 1.9, fill=LIGHT)
    add_text(sl, detail, x + 0.15, y + 0.75, 3.45, 1.6,
             font_size=13, color=DARK, align=PP_ALIGN.LEFT)

# Arrows between phases
for i in [0, 1, 3, 4]:
    ax = 4.15 + (i % 3) * 4.2
    ay = 1.4 + (i // 3) * 2.7 + 1.2
    add_text(sl, "→", ax - 0.1, ay, 0.5, 0.4, font_size=22, bold=True, color=NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATA SOURCES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Data Sources", "Five Datasets from HDX / DHS Program")
footer(sl, 5)

sources = [
    ("DHS / MICS Surveys", "Stunting, wasting, maternal health, WASH\n5 waves: 1991, 1998, 2004, 2011, 2018", BLUE),
    ("WHO Health Indicators", "Malaria prevalence, child anaemia, health insurance\nMatched to survey years", NAVY),
    ("Cameroon Admin Boundaries", "GeoJSON for choropleth visualisation\n10 administrative regions", ORANGE),
    ("Health Facility Locations", "OpenStreetMap facility counts per region\nCurrent coverage", RGBColor(0, 112, 192)),
    ("Population Density Maps", "WorldPop 2020 high-resolution rasters\nRegional population weighting", GREEN),
]

for i, (name, detail, col) in enumerate(sources):
    x = 0.35 + (i % 2) * 6.5 if i < 4 else 3.5
    y = 1.4 + (i // 2) * 1.85
    if i == 4:
        y = 5.2
    add_rect(sl, x, y, 6.1, 1.65, fill=LIGHT)
    add_rect(sl, x, y, 0.25, 1.65, fill=col)
    add_text(sl, name, x + 0.35, y + 0.12, 5.6, 0.4,
             font_size=15, bold=True, color=NAVY)
    add_text(sl, detail, x + 0.35, y + 0.52, 5.6, 1.0,
             font_size=12, color=GREY)

add_text(sl, "→  All sources publicly available from HDX (data.humdata.org) — no DUA required",
         0.4, 7.0, 12.5, 0.25, font_size=11, italic=True, color=GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DATASET OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Dataset Overview", "Analysis-Ready Wide Table: 50 rows × 20 columns")
footer(sl, 6)

# Stats boxes
stats = [
    ("50", "Region-Year\nObservations"),
    ("10", "Cameroon\nRegions"),
    ("5", "Survey\nWaves"),
    ("15", "Raw\nFeatures"),
    ("4", "Engineered\nFeatures"),
    ("1", "Target\n(stunting_rate)"),
]
for i, (val, lbl) in enumerate(stats):
    x = 0.4 + i * 2.1
    add_rect(sl, x, 1.3, 1.85, 1.3, fill=BLUE if i < 3 else ORANGE)
    add_text(sl, val, x, 1.38, 1.85, 0.75,
             font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, x, 2.05, 1.85, 0.5,
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Feature Groups", 0.4, 2.85, 12.0, 0.35, font_size=15, bold=True, color=NAVY)

feat_groups = [
    ("Maternal Education", "women_secondary_plus_pct\nwomen_literate_pct"),
    ("WASH", "improved_water_pct\nimproved_sanitation_pct"),
    ("Healthcare Access", "antenatal_4plus_pct\nskilled_birth_attendance_pct\nhealth_facility_delivery_pct\nhealth_insurance_any_pct"),
    ("Disease Burden", "malaria_prevalence_pct\nchild_anemia_pct"),
    ("Nutrition Siblings", "wasting_pct\nunderweight_pct"),
    ("Engineered", "wash_composite\nmaternal_health_score\neducation_score\ndisease_burden"),
]

for i, (grp, feats) in enumerate(feat_groups):
    x = 0.35 + (i % 3) * 4.25
    y = 3.3 + (i // 3) * 1.8
    add_rect(sl, x, y, 4.0, 0.38, fill=NAVY if i < 3 else ORANGE)
    add_text(sl, grp, x + 0.1, y + 0.04, 3.85, 0.3,
             font_size=13, bold=True, color=WHITE)
    add_text(sl, feats, x + 0.1, y + 0.42, 3.85, 1.25,
             font_size=11, color=GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EDA: TARGET DISTRIBUTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Exploratory Data Analysis", "Target Distribution — Stunting Rate")
footer(sl, 7)

add_image_or_placeholder(sl, "01_target_distribution.png",
                          0.4, 1.3, 6.5, 4.9,
                          "Figure 1: Distribution of Stunting Rate (all 50 region-year obs.)")

add_text(sl, "Key Observations", 7.2, 1.3, 5.6, 0.4, font_size=16, bold=True, color=NAVY)
add_bullet_block(sl, [
    "Range: 15.4% (Littoral 2018) → 45.7% (Far North peak)",
    "Mean ≈ 30% — at the WHO 'high' risk threshold",
    "~40% of observations exceed the 30% WHO action level",
    "Slight right skew — extreme northern observations pull the tail",
    "Clear bimodal tendency: northern cluster (35-40%) vs southern cluster (20-25%)",
], 7.2, 1.75, 5.7, 4.4, font_size=13)

add_rect(sl, 7.2, 5.4, 5.7, 0.8, fill=LIGHT)
add_text(sl, "WHO Risk Thresholds:  Low < 20%  |  Medium 20–29%  |  High 30–39%  |  Critical ≥ 40%",
         7.3, 5.5, 5.5, 0.6, font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — EDA: REGIONAL PATTERNS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Exploratory Data Analysis", "Regional Patterns — North-South Gradient")
footer(sl, 8)

add_image_or_placeholder(sl, "02_region_boxplot.png",
                          0.4, 1.3, 7.8, 4.9,
                          "Figure 2: Stunting Rate by Region (median-sorted boxplot, 1991–2018)")

add_text(sl, "North-South Gradient", 8.5, 1.3, 4.5, 0.4, font_size=16, bold=True, color=NAVY)

regions_risk = [
    ("Far North", "37.9%", RED),
    ("North", "37.2%", RED),
    ("Adamawa", "36.5%", ORANGE),
    ("East", "32.1%", ORANGE),
    ("West", "27.4%", BLUE),
    ("Centre", "21.8%", GREEN),
    ("Littoral", "15.4%", GREEN),
]
for i, (reg, val, col) in enumerate(regions_risk):
    y = 1.85 + i * 0.55
    add_rect(sl, 8.5, y, 3.5, 0.45, fill=LIGHT)
    add_rect(sl, 8.5, y, 0.18, 0.45, fill=col)
    add_text(sl, reg, 8.75, y + 0.06, 2.2, 0.35, font_size=12, color=DARK)
    add_text(sl, val, 11.6, y + 0.06, 0.8, 0.35,
             font_size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)

add_text(sl, "⬆ Critical / High risk", 8.5, 5.85, 4.5, 0.3,
         font_size=11, italic=True, color=RED)
add_text(sl, "⬇ Low risk", 8.5, 6.15, 4.5, 0.3,
         font_size=11, italic=True, color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EDA: CORRELATION HEATMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Exploratory Data Analysis", "Correlation Heatmap — Feature Relationships")
footer(sl, 9)

add_image_or_placeholder(sl, "03_correlation_heatmap.png",
                          0.4, 1.3, 7.5, 5.0,
                          "Figure 3: Spearman Correlation Heatmap (features vs stunting rate)")

add_text(sl, "Top Correlations with Stunting", 8.2, 1.3, 4.8, 0.4,
         font_size=15, bold=True, color=NAVY)

corrs = [
    ("health_facility_delivery_pct", "ρ = −0.78", RED),
    ("SES proxy", "ρ = −0.71", RED),
    ("women_secondary_plus_pct", "ρ = −0.54", ORANGE),
    ("wash_composite", "ρ = −0.50", ORANGE),
    ("malaria_prevalence_pct", "ρ = +0.36", BLUE),
    ("child_anemia_pct", "ρ = +0.29", BLUE),
]
for i, (feat, val, col) in enumerate(corrs):
    y = 1.8 + i * 0.65
    add_rect(sl, 8.2, y, 4.8, 0.55, fill=LIGHT)
    add_text(sl, feat, 8.35, y + 0.08, 3.2, 0.42, font_size=11, color=DARK)
    add_text(sl, val, 11.45, y + 0.08, 1.45, 0.42,
             font_size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)

add_text(sl, "→ Healthcare access is the dominant correlate\n   Education & WASH follow closely",
         8.2, 6.0, 4.8, 0.65, font_size=12, italic=True, color=NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DATA PREPARATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Data Preparation", "Cleaning → Imputation → Feature Engineering")
footer(sl, 10)

steps = [
    ("1  Type Coercion", "Cast all % columns to float64; year to int"),
    ("2  De-duplication", "Resolve duplicate region-year-indicator rows from overlapping survey rounds"),
    ("3  Name Normalisation", "Lookup table maps 25+ DHS region label variants to 10 canonical names"),
    ("4  Mega-region Expansion", "1991/1998 large regions broadcast to all constituent sub-regions (flagged as imputed)"),
    ("5  Missing Value Imputation", "Level 1: region-year median  →  Level 2: global median for remaining gaps"),
]

for i, (step, detail) in enumerate(steps):
    y = 1.35 + i * 0.95
    add_rect(sl, 0.35, y, 0.7, 0.75, fill=ORANGE)
    add_text(sl, str(i + 1), 0.35, y + 0.1, 0.7, 0.55,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.15, y, 11.8, 0.75, fill=LIGHT)
    add_text(sl, step, 1.3, y + 0.04, 3.5, 0.35, font_size=13, bold=True, color=NAVY)
    add_text(sl, detail, 1.3, y + 0.38, 11.5, 0.35, font_size=12, color=DARK)

add_text(sl, "Engineered Features", 0.35, 6.25, 12.0, 0.35, font_size=15, bold=True, color=NAVY)

eng = [
    ("wash_composite", "Mean(improved_water, improved_sanitation)"),
    ("maternal_health_score", "Mean(antenatal_4plus, skilled_birth, facility_delivery)"),
    ("education_score", "Mean(literate_pct, secondary_plus_pct)"),
    ("disease_burden", "Mean(malaria_prevalence, child_anemia)"),
]
for i, (name, formula) in enumerate(eng):
    x = 0.35 + i * 3.1
    add_rect(sl, x, 6.62, 2.9, 0.55, fill=BLUE)
    add_text(sl, name, x + 0.1, 6.65, 2.75, 0.22,
             font_size=10, bold=True, color=WHITE)
    add_text(sl, formula, x + 0.1, 6.88, 2.75, 0.3,
             font_size=9, italic=True, color=LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MODELLING: MODEL ZOO
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Regression Modelling", "Ten Models Evaluated with 5-Fold Cross-Validation")
footer(sl, 11)

models = [
    ("XGBoost ⭐", "2.87 ± 0.70", GREEN),
    ("Random Forest", "3.16 ± 0.93", GREEN),
    ("Gradient Boosting", "3.19 ± 0.74", GREEN),
    ("LightGBM", "3.45 ± 0.88", BLUE),
    ("Decision Tree", "4.12 ± 1.21", BLUE),
    ("Ridge Regression", "4.67 ± 0.95", ORANGE),
    ("Lasso Regression", "4.89 ± 1.03", ORANGE),
    ("Linear Regression", "5.14 ± 1.18", ORANGE),
    ("K-Nearest Neighbours", "5.78 ± 1.44", RED),
    ("Baseline (mean) ⬇", "7.17 ± 1.45", RED),
]

add_text(sl, "Model", 0.35, 1.3, 5.5, 0.35, font_size=14, bold=True, color=NAVY)
add_text(sl, "CV RMSE (pp)", 6.0, 1.3, 3.5, 0.35, font_size=14, bold=True, color=NAVY)
add_text(sl, "Bar", 9.7, 1.3, 3.3, 0.35, font_size=14, bold=True, color=NAVY)

for i, (name, rmse, col) in enumerate(models):
    y = 1.75 + i * 0.52
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(sl, 0.35, y, 5.5, 0.48, fill=bg)
    add_text(sl, name, 0.45, y + 0.06, 5.3, 0.38, font_size=12,
             bold=(i == 0), color=DARK if i != 0 else NAVY)
    add_rect(sl, 6.0, y, 3.5, 0.48, fill=bg)
    add_text(sl, rmse, 6.1, y + 0.06, 3.3, 0.38, font_size=13, bold=(i == 0),
             color=col, align=PP_ALIGN.CENTER)
    # Bar chart representation
    bar_w = (float(rmse.split()[0]) / 7.17) * 3.0
    add_rect(sl, 9.7, y + 0.1, bar_w, 0.3, fill=col)

add_text(sl, "XGBoost achieves 60% RMSE improvement over baseline  (p = 0.003)",
         0.35, 7.0, 12.5, 0.3, font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — MODEL LEADERBOARD CHART
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Regression Results", "Model Leaderboard & Predicted vs Actual")
footer(sl, 12)

add_image_or_placeholder(sl, "06_model_leaderboard.png",
                          0.35, 1.3, 6.4, 4.8,
                          "Figure 5: 5-Fold CV RMSE — All 10 Models")

add_image_or_placeholder(sl, "09_predicted_vs_actual.png",
                          7.0, 1.3, 6.0, 4.8,
                          "Figure 12: Predicted vs Actual (Hold-out Test Set, XGBoost)")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — FEATURE IMPORTANCE / SHAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Feature Importance & SHAP Analysis", "What drives the XGBoost predictions?")
footer(sl, 13)

add_image_or_placeholder(sl, "07_feature_importance.png",
                          0.35, 1.3, 6.2, 4.8,
                          "Figure 6: Top 15 Features — XGBoost Gain Importance")

add_image_or_placeholder(sl, "08_shap_beeswarm.png",
                          6.8, 1.3, 6.2, 4.8,
                          "Figure 7: SHAP Beeswarm — Observation-level Attribution")

add_text(sl, "underweight_pct dominates (~0.78 SHAP)  |  "
             "Healthcare + Education + WASH ≈ 60% of remaining importance",
         0.35, 6.35, 12.5, 0.35, font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — HYPOTHESIS TESTING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Hypothesis Testing", "All 6 Null Hypotheses Rejected at p < 0.05")
footer(sl, 14)

hyps = [
    ("H1", "Maternal Education", "ρ = −0.538", "p = 5.66×10⁻⁵", "✓"),
    ("H2", "WASH Composite", "ρ = −0.501", "p = 2.12×10⁻⁴", "✓"),
    ("H3", "SES Proxy", "ρ = −0.709", "p = 8.35×10⁻⁹", "✓"),
    ("H4", "Healthcare Access", "ρ = −0.778", "p = 2.84×10⁻¹¹", "✓"),
    ("H5", "Malaria Prevalence", "ρ = +0.358", "p = 1.06×10⁻²", "✓"),
    ("H6", "ML Beats Baseline", "−60% RMSE", "p = 0.003", "✓"),
]

strength_cols = [ORANGE, ORANGE, BLUE, RED, BLUE, GREEN]

for i, (h, driver, stat, pval, decision) in enumerate(hyps):
    x = 0.35 + (i % 3) * 4.25
    y = 1.4 + (i // 3) * 2.7
    add_rect(sl, x, y, 3.95, 0.55, fill=strength_cols[i])
    add_text(sl, f"{h}: {driver}", x + 0.15, y + 0.08, 3.7, 0.42,
             font_size=15, bold=True, color=WHITE)
    add_rect(sl, x, y + 0.55, 3.95, 1.85, fill=LIGHT)
    add_text(sl, f"Statistic: {stat}", x + 0.15, y + 0.7, 3.7, 0.4,
             font_size=13, color=DARK)
    add_text(sl, pval, x + 0.15, y + 1.1, 3.7, 0.4, font_size=13, color=GREY)
    add_text(sl, f"Null rejected {decision}", x + 0.15, y + 1.5, 3.7, 0.4,
             font_size=14, bold=True, color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLASSIFICATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "WHO Risk Band Classification", "Assigning Regions to Low / Medium / High / Critical")
footer(sl, 15)

add_image_or_placeholder(sl, "11_confusion_matrix.png",
                          0.35, 1.3, 6.5, 5.0,
                          "Figure 8: Confusion Matrix — Random Forest Classifier (80% accuracy)")

add_text(sl, "Risk Band Thresholds", 7.2, 1.3, 5.8, 0.4, font_size=16, bold=True, color=NAVY)
bands = [
    ("Critical", "≥ 40%", RED),
    ("High", "30 – 39%", ORANGE),
    ("Medium", "20 – 29%", BLUE),
    ("Low", "< 20%", GREEN),
]
for i, (band, threshold, col) in enumerate(bands):
    y = 1.85 + i * 0.7
    add_rect(sl, 7.2, y, 5.8, 0.6, fill=col)
    add_text(sl, f"{band}  →  Stunting {threshold}", 7.35, y + 0.1, 5.5, 0.4,
             font_size=15, bold=True, color=WHITE)

add_text(sl, "Model Performance", 7.2, 4.75, 5.8, 0.4, font_size=16, bold=True, color=NAVY)
add_bullet_block(sl, [
    "Random Forest: 80% accuracy (best)",
    "Logistic Regression: 73% accuracy",
    "Errors confined to adjacent bands only",
    "No cross-two-band misclassifications",
], 7.2, 5.2, 5.8, 2.0, font_size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CLUSTERING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Clustering Analysis", "K-Means Discovers Two Distinct Region Archetypes")
footer(sl, 16)

add_image_or_placeholder(sl, "12_elbow.png",
                          0.35, 1.3, 5.5, 3.0,
                          "Figure 9: Elbow Curve — Choosing k=2")
add_image_or_placeholder(sl, "13_clusters_pca.png",
                          6.1, 1.3, 7.0, 3.0,
                          "Figure 10: Clusters in PCA Space (k=2, silhouette=0.52)")

# Cluster profiles
add_text(sl, "Cluster Profiles", 0.35, 4.55, 12.7, 0.35, font_size=15, bold=True, color=NAVY)

cl_data = [
    ("Cluster 0: Northern-rural (Critical Risk)", RED,
     "Far North, North, Adamawa, East\nWomen Secondary+: 1.9%  |  Literate: 43.6%  |  Stunting: 37.0%"),
    ("Cluster 1: Urban-affluent (Low Risk)", GREEN,
     "Centre, Littoral, South, West, North-West, South-West\nWomen Secondary+: 10.3%  |  Literate: 87.2%  |  Stunting: 23.6%"),
]
for i, (name, col, detail) in enumerate(cl_data):
    x = 0.35 + i * 6.5
    add_rect(sl, x, 5.0, 6.2, 0.45, fill=col)
    add_text(sl, name, x + 0.15, 5.05, 5.9, 0.38,
             font_size=13, bold=True, color=WHITE)
    add_rect(sl, x, 5.45, 6.2, 1.0, fill=LIGHT)
    add_text(sl, detail, x + 0.15, 5.52, 5.9, 0.9, font_size=12, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — FORECASTING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Forecasting to 2026 / 2028", "Linear Time-Trend Extrapolation per Region")
footer(sl, 17)

add_image_or_placeholder(sl, "14_forecasts_top5.png",
                          0.35, 1.3, 7.5, 5.1,
                          "Figure 11: Stunting Rate Forecasts — Top 5 Hotspot Regions")

add_text(sl, "Forecast Summary", 8.1, 1.3, 4.9, 0.4, font_size=16, bold=True, color=NAVY)

forecast_data = [
    ("Far North", "37.9%", "37.8%", "37.6%", "−0.09 pp/yr", RED),
    ("North", "37.2%", "37.1%", "36.9%", "−0.10 pp/yr", RED),
    ("Adamawa", "36.5%", "36.3%", "36.1%", "−0.15 pp/yr", ORANGE),
    ("Littoral", "15.4%", "15.3%", "15.0%", "−0.29 pp/yr", GREEN),
]

for i, (reg, now, f26, f28, slope, col) in enumerate(forecast_data):
    y = 1.85 + i * 0.9
    add_rect(sl, 8.1, y, 4.9, 0.8, fill=LIGHT)
    add_rect(sl, 8.1, y, 0.2, 0.8, fill=col)
    add_text(sl, reg, 8.4, y + 0.05, 2.5, 0.35, font_size=13, bold=True, color=DARK)
    add_text(sl, f"Now: {now}  →  2026: {f26}  →  2028: {f28}",
             8.4, y + 0.42, 4.5, 0.32, font_size=10, color=GREY)

add_rect(sl, 8.1, 5.6, 4.9, 0.75, fill=RED)
add_text(sl, "Far North remains in 'high' WHO band through 2028+\nat current improvement rate of −0.09 pp/year",
         8.2, 5.65, 4.7, 0.65, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — HOTSPOT RANKINGS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Hotspot Rankings", "Ranked Targeting Intelligence for Ministry of Public Health")
footer(sl, 18)

add_image_or_placeholder(sl, "10_hotspots_top15.png",
                          0.35, 1.3, 6.5, 5.0,
                          "Figure 13: Top 15 Predicted Stunting Hotspots (2018, XGBoost scores)")

add_text(sl, "Top Hotspot Regions (2018)", 7.2, 1.3, 5.8, 0.4,
         font_size=15, bold=True, color=NAVY)

hotspots = [
    ("1", "Far North", "37.9%", "Critical", RED),
    ("2", "North", "37.2%", "Critical", RED),
    ("3", "Adamawa", "36.5%", "High", ORANGE),
    ("4", "East", "32.1%", "High", ORANGE),
    ("5", "North West", "29.8%", "Medium", BLUE),
    ("6", "South West", "28.3%", "Medium", BLUE),
    ("7", "West", "27.4%", "Medium", BLUE),
]
for i, (rank, region, pred, band, col) in enumerate(hotspots):
    y = 1.82 + i * 0.59
    add_rect(sl, 7.2, y, 0.4, 0.5, fill=col)
    add_text(sl, rank, 7.2, y + 0.06, 0.4, 0.38,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 7.65, y, 5.3, 0.5, fill=LIGHT)
    add_text(sl, region, 7.78, y + 0.06, 2.5, 0.38, font_size=12, color=DARK)
    add_text(sl, pred, 10.35, y + 0.06, 1.1, 0.38,
             font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, band, 11.5, y + 0.06, 1.4, 0.38,
             font_size=11, color=col, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Deployment", "Client-Side Web Application + Offline Artefacts")
footer(sl, 19)

add_text(sl, "Web Application — 8 Interactive Pages", 0.35, 1.3, 8.5, 0.4,
         font_size=16, bold=True, color=NAVY)

pages = [
    "Home — KPIs, H1-H6 hypothesis strip, top-10 hotspots",
    "Hotspots — Full regional stunting ranking table",
    "What-If Predictor — Client-side inference (no server needed)",
    "Regression — Interactive model leaderboard",
    "Classification — Confusion matrix + per-class F1",
    "Clustering — Cluster profiles + silhouette",
    "Forecasts — Per-region trends + 2026/2028 projections",
    "Hypotheses — H1-H6 cards with test statistics",
]
add_bullet_block(sl, pages, 0.35, 1.8, 8.3, 5.0, font_size=12)

add_text(sl, "Tech Stack", 9.0, 1.3, 4.0, 0.4, font_size=16, bold=True, color=NAVY)
tech = [
    ("Next.js 16 + React 19", BLUE),
    ("TypeScript + Tailwind CSS 4", NAVY),
    ("Redux Toolkit (state)", ORANGE),
    ("Plotly.js (charts)", GREEN),
    ("predict.ts (client inference)", RED),
    ("predictor.json (model export)", GREY),
]
for i, (t, col) in enumerate(tech):
    y = 1.82 + i * 0.65
    add_rect(sl, 9.0, y, 4.0, 0.55, fill=col)
    add_text(sl, t, 9.1, y + 0.08, 3.8, 0.38,
             font_size=13, bold=True, color=WHITE)

add_text(sl, "→ Offline-capable: no server call for predictions",
         9.0, 6.15, 4.0, 0.3, font_size=11, italic=True, color=NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — KEY FINDINGS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Key Findings", "Five Evidence-Based Conclusions")
footer(sl, 20)

findings = [
    ("North-South Gradient is Real", NAVY,
     "Far North (37.9%), North (37.2%), Adamawa (36.5%) consistently dominate the hotspot list "
     "across all survey waves. The gap with southern regions exceeds 22 pp."),
    ("Healthcare Access Dominates", RED,
     "The strongest correlate of stunting is healthcare access (ρ = −0.78, p < 10⁻¹⁰). "
     "Improving facility delivery and antenatal care is the highest-leverage single intervention."),
    ("ML Beats Baseline by 60%", GREEN,
     "XGBoost achieves CV RMSE = 2.87 pp vs baseline 7.17 pp (p = 0.003), validating "
     "machine-learning for subnational malnutrition prediction."),
    ("Two Distinct Clusters", ORANGE,
     "K-Means (k=2, silhouette=0.52) finds a critical northern cluster (37% stunting) "
     "and a low-risk southern cluster (24%). Bundle interventions are needed in the north."),
    ("Insufficient Improvement Rate", BLUE,
     "Linear forecasts show northern regions remain in WHO 'high'/'critical' band through 2028+. "
     "Accelerated, multi-sectoral action is urgently needed."),
]

for i, (title, col, detail) in enumerate(findings):
    x = 0.35 + (i % 2) * 6.5 if i < 4 else 3.5
    y = 1.4 + (i // 2) * 2.35
    if i == 4:
        y = 5.6
        x = 0.35
    add_rect(sl, x, y, 6.1 if i < 4 else 12.5, 0.5, fill=col)
    add_text(sl, f"{i+1}. {title}", x + 0.15, y + 0.07, 5.8 if i < 4 else 12.2, 0.38,
             font_size=15, bold=True, color=WHITE)
    add_rect(sl, x, y + 0.5, 6.1 if i < 4 else 12.5, 1.55, fill=LIGHT)
    add_text(sl, detail, x + 0.15, y + 0.62, 5.8 if i < 4 else 12.2, 1.3,
             font_size=12, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — CONCLUSIONS & RECOMMENDATIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Conclusions & Recommendations", "Action Plan for the Ministry of Public Health")
footer(sl, 21)

add_text(sl, "Policy Recommendations", 0.35, 1.3, 12.5, 0.4, font_size=16, bold=True, color=NAVY)

recs = [
    ("Priority Targeting", ORANGE,
     "Concentrate first-wave investments on Far North, North, and Adamawa — "
     "all three appear in top-3 across every analytical lens."),
    ("Bundle Interventions", BLUE,
     "Cluster analysis shows high-risk regions suffer simultaneous deficits in education, WASH, "
     "and healthcare. Single-sector approaches are insufficient."),
    ("Monitor Improvement Velocity", NAVY,
     "Regions with the slowest improvement slopes should receive disproportionately higher "
     "investment to close the north-south gap."),
    ("Use the What-If Predictor", GREEN,
     "Train district health officers to simulate the impact of coverage improvements "
     "before committing intervention budgets."),
    ("Update with 2024/2025 DHS Data", RED,
     "Re-run the pipeline when the next DHS wave is published to validate forecasts "
     "and re-rank hotspots."),
]

for i, (title, col, detail) in enumerate(recs):
    y = 1.82 + i * 0.98
    add_rect(sl, 0.35, y, 0.35, 0.85, fill=col)
    add_text(sl, str(i + 1), 0.35, y + 0.18, 0.35, 0.5,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 0.75, y, 12.2, 0.85, fill=LIGHT)
    add_text(sl, title, 0.9, y + 0.05, 3.5, 0.38, font_size=13, bold=True, color=col)
    add_text(sl, detail, 0.9, y + 0.45, 11.8, 0.38, font_size=12, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — REFERENCES / THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 2.5, 13.33, 2.5, fill=BLUE)
add_rect(sl, 0, 2.46, 13.33, 0.08, fill=ORANGE)
add_rect(sl, 0, 4.96, 13.33, 0.08, fill=ORANGE)

add_text(sl, "Thank You", 1.0, 0.4, 11.0, 1.0,
         font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri Light")
add_text(sl, "Questions & Discussion", 1.0, 1.35, 11.0, 0.5,
         font_size=20, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(sl, "Selected References", 1.0, 2.65, 11.0, 0.38,
         font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

refs_short = [
    "[1] DHS Program. (2018). Cameroon DHS 2018. ICF International.",
    "[2] Chen & Guestrin. (2016). XGBoost. KDD '16.",
    "[3] Chapman et al. (2000). CRISP-DM 1.0. SPSS Inc.",
    "[4] WHO. (2014). Global Nutrition Targets 2025. WHO.",
    "[5] Humanitarian Data Exchange. (2024). data.humdata.org",
]
for i, ref in enumerate(refs_short):
    y = 3.1 + i * 0.35
    add_text(sl, ref, 1.5, y, 10.0, 0.32, font_size=10, italic=True, color=LIGHT)

add_text(sl, "SEPO PERRY-BRADLEY DINGA  |  CT23A145  |  CEC 420  |  University of Buea  |  June 2026",
         0.5, 5.2, 12.3, 0.35, font_size=13, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "GitHub: Cameroon Malnutrition Atlas  |  Built with Python 3.11 + Next.js 16",
         0.5, 5.6, 12.3, 0.35, font_size=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

# footer
add_rect(sl, 0, 7.2, 13.33, 0.3, fill=RGBColor(0, 31, 63))
add_text(sl, "22 / 22", 12.0, 7.22, 1.0, 0.25, font_size=9, color=LIGHT, align=PP_ALIGN.RIGHT)

# ── save ───────────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Presentation saved -> {OUT}")
